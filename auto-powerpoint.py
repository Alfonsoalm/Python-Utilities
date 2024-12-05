import os
from pptx import Presentation
from pptx.util import Inches


# === Funcion para crear o añadir a un powerpoint imagenes de una carpeta ===
def create_presentation_with_images(folder_path, ppt_path):
    try:
        # Verificar si el archivo ya existe
        if os.path.exists(ppt_path):
            print(f"El archivo '{ppt_path}' existe. Se abrirá para editar.")
            presentation = Presentation(ppt_path)
        else:
            print(f"El archivo '{ppt_path}' no existe. Se creará uno nuevo.")
            presentation = Presentation()

        # Obtener todas las imágenes de la carpeta
        image_files = sorted(
            [os.path.join(folder_path, f) for f in os.listdir(folder_path) if f.lower().endswith(('png', 'jpg', 'jpeg'))]
        )

        if not image_files:
            print("No se encontraron imágenes en la carpeta.")
            return

        print(f"Imágenes encontradas: {image_files}")

        # Agregar imágenes a las diapositivas
        image_index = 0
        for slide in presentation.slides:
            if image_index >= len(image_files):
                break  # No quedan más imágenes

            # Agregar la primera imagen
            left = Inches(0.5)  # Posición izquierda
            top = Inches(3)    # Posición superior
            slide.shapes.add_picture(image_files[image_index], left, top, width=Inches(4))
            image_index += 1

            # Agregar la segunda imagen si hay más
            if image_index < len(image_files):
                left = Inches(5)  # Posición izquierda para la segunda imagen
                top = Inches(3)   # Misma posición superior
                slide.shapes.add_picture(image_files[image_index], left, top, width=Inches(4))
                image_index += 1

        # Si quedan imágenes sin diapositiva, crear nuevas diapositivas
        while image_index < len(image_files):
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Diapositiva en blanco

            # Agregar la primera imagen
            left = Inches(0.5)
            top = Inches(1)
            slide.shapes.add_picture(image_files[image_index], left, top, width=Inches(4))
            image_index += 1

            # Agregar la segunda imagen si hay más
            if image_index < len(image_files):
                left = Inches(5)
                top = Inches(1)
                slide.shapes.add_picture(image_files[image_index], left, top, width=Inches(4))
                image_index += 1

        # Guardar la presentación
        presentation.save(ppt_path)
        print(f"Presentación guardada exitosamente: {ppt_path}")

    except Exception as e:
        print(f"Error al procesar la presentación: {e}")

if __name__ == "__main__":
    # Ruta de la carpeta con las imágenes
    folder_path = "imagenes"  # Cambia esto por la ruta real

    # Nombre del archivo de salida
    ppt_path = "documentos_generados/pptx/presentacion_imagenes.pptx"

    # Crear la presentación
    create_presentation_with_images(folder_path, ppt_path)
